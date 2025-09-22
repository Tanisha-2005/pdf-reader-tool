import os
import json
import io
import fitz
import pkg_resources
from tkinter import Tk, Frame, Canvas, Text, Scrollbar, Button, Entry, Label, Listbox, X, Toplevel, simpledialog, messagebox, BOTH, END, LEFT, RIGHT, Y, W, WORD
from tkinter import filedialog
from PIL import Image, ImageTk, ImageDraw
from tkinter import SINGLE, DISABLED

from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lex_rank import LexRankSummarizer

from transformers import pipeline
from docx import Document
from pptx import Presentation
from ebooklib import epub, ITEM_DOCUMENT
from bs4 import BeautifulSoup
import html2text
from PyPDF2 import PdfReader, PdfWriter
from cryptography.fernet import Fernet
from cryptography.hazmat.primitives import hashes
from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
import base64
import mammoth
from tkhtmlview import HTMLLabel
import logging

# Configure logging for detailed debugging
logging.basicConfig(level=logging.DEBUG, format="%(asctime)s - %(levelname)s - %(message)s")

def protect_pdf(input_path, output_path, password):
    """Encrypt a PDF file using PyPDF2 with AES-128 and common permissions."""
    try:
        reader = PdfReader(input_path)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        # Use AES-128 encryption with permissions for printing, copying, and annotating
        writer.encrypt(
            user_password=password,
            owner_password=password,
            permissions_flag=0b111100110100,  # Allow print, copy, annotate
            use_128bit=True
        )

        with open(output_path, "wb") as f:
            writer.write(f)
        logging.debug(f"Encrypted PDF created: {output_path} with password (length: {len(password)})")
        return True
    except Exception as e:
        logging.error(f"Failed to encrypt PDF {input_path}: {str(e)}")
        raise ValueError(f"Failed to encrypt PDF: {str(e)}")

class SmartDocReader:
    def __init__(self, root):

        self.root = root
        self.current_page = 0   
        self.bookmarks = []     
        self.total_pages = 10  # Example, you need to set this when opening a file

        self.root = root
        self.root.title("📄 Smart Document Reader with Summarizer")
        self.root.geometry("1200x850")
        self.root.configure(bg="#1e1e1e")

       
        self.page_widgets = []
        self.tk_images = []
        self.drawing = False
        self.canvas = None
        self.inner_frame = None
        self.canvas_window = None
        self.file_path = ""
        self.text_box = None
        self.text_entry = None
        self.doc_text = ""
        self.current_text_widget = None
        self.protected_files = self.load_protected_files()
        self.page_count_label = None
        self.pdf_doc = None  # Initialize class-level PDF document storage
        self.eraser_mode = False
        self.draw_highlights = {}
        self.search_highlights = {}

        self.build_ui()
        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)

    def on_closing(self):
        """Close pdf_doc and destroy window."""
        if self.pdf_doc:
            self.pdf_doc.close()
            self.pdf_doc = None
        self.root.destroy()

    def load_protected_files(self):
        """Load protected files from a JSON file."""
        try:
            with open("protected_files.json", "r") as f:
                data = json.load(f)
                logging.debug(f"Loaded protected files: {list(data.keys())}")
                return data
        except (FileNotFoundError, json.JSONDecodeError):
            logging.debug("No protected_files.json found or invalid, returning empty dict")
            return {}

    def save_protected_files(self):
        """Save protected files to a JSON file."""
        try:
            with open("protected_files.json", "w") as f:
                json.dump(self.protected_files, f, indent=4)
            logging.debug("Saved protected_files.json")
        except Exception as e:
            logging.error(f"Failed to save protected files: {str(e)}")
            messagebox.showerror("Error", f"Failed to save protected files: {str(e)}")

    def hash_password(self, password):
        """Hash a password using PBKDF2."""
        kdf = PBKDF2HMAC(
            algorithm=hashes.SHA256(),
            length=32,
            salt=b'password_salt_',
            iterations=100000,
        )
        key = base64.urlsafe_b64encode(kdf.derive(password.encode('utf-8')))
        return key.decode('utf-8')

    def prompt_password(self, file_path):
        """Prompt for password and handle PDF/non-PDF cases."""
        file_path = os.path.normpath(file_path)
        stored_password = self.protected_files.get(file_path)
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf" and stored_password:
            logging.debug(f"Using stored password for PDF {file_path} (length: {len(stored_password)})")
            if self.test_pdf_encryption_fitz(file_path, stored_password):
                logging.debug(f"Password verified with PyMuPDF for {file_path}")
                return stored_password
            logging.warning(f"Stored password failed for {file_path}, prompting user")

        password = simpledialog.askstring("Password", f"Enter password for {os.path.basename(file_path)}:", show="*")
        if password is None:
            logging.debug(f"Password prompt cancelled for {file_path}")
            return None
        password = password.strip()
        logging.debug(f"Entered password for {file_path} (length: {len(password)})")

        if ext == ".pdf":
            if self.test_pdf_encryption_fitz(file_path, password):
                logging.debug(f"Password verified with PyMuPDF for {file_path}")
                return password
            logging.error(f"Password verification failed with PyMuPDF for {file_path}")
            messagebox.showerror("Error", "Incorrect password.")
            return None
        else:
            if file_path in self.protected_files:
                hashed_password = self.hash_password(password)
                if hashed_password != self.protected_files[file_path]:
                    logging.error(f"Incorrect password for non-PDF {file_path}")
                    messagebox.showerror("Error", "Incorrect password.")
                    return None
            logging.debug(f"Password verified for non-PDF {file_path}")
            return password

    def test_pdf_encryption_pypdf2(self, file_path, password):
        """Test if a PDF can be opened with PyPDF2."""
        try:
            reader = PdfReader(file_path)
            if reader.is_encrypted:
                result = reader.decrypt(password)
                if result == 0:
                    logging.error(f"PyPDF2 decryption failed for {file_path}")
                    return False
                reader.pages[0]  # Access a page to ensure decryption worked
            logging.debug(f"PyPDF2 decryption succeeded for {file_path}")
            return True
        except Exception as e:
            logging.error(f"PyPDF2 test failed for {file_path}: {str(e)}")
            return False

    def test_pdf_encryption_fitz(self, file_path, password):
        """Test if a PDF can be opened with PyMuPDF."""
        try:
            doc = fitz.open(file_path)
            if doc.is_encrypted:
                if not doc.authenticate(password):
                    doc.close()
                    logging.error(f"PyMuPDF authentication failed for {file_path}")
                    return False
            doc.close()
            logging.debug(f"PyMuPDF authentication succeeded for {file_path}")
            return True
        except Exception as e:
            logging.error(f"PyMuPDF test failed for {file_path}: {str(e)}")
            return False

    def build_ui(self):
        """Build the user interface."""
        top_frame1 = Frame(self.root, bg="#1e1e1e")
        top_frame1.pack(fill=X, padx=10, pady=(10, 2))

        top_frame2 = Frame(self.root, bg="#1e1e1e")
        top_frame2.pack(fill=X, padx=10, pady=(2, 10))

        Button(top_frame1, text="Open File", command=self.select_file).pack(side=LEFT, padx=5)
        Button(top_frame1, text="Summarize 🧠", command=self.summarize_text).pack(side=LEFT, padx=5)
        self.text_entry = Entry(top_frame1, width=30)
        self.text_entry.pack(side=LEFT, padx=5)
        self.text_entry.bind("<Return>", lambda event: self.search_keyword())
        Button(top_frame1, text="Search", command=self.search_keyword).pack(side=LEFT, padx=5)
        Button(top_frame2, text="Clear Highlights", command=self.clear_highlights).pack(side=LEFT, padx=5)
        Button(top_frame2, text="Copy Text 📋", command=self.copy_text_popup).pack(side=LEFT, padx=5)
        Button(top_frame2, text="Draw Mode ✏", command=self.toggle_draw_mode).pack(side=LEFT, padx=5)
        Button(top_frame2, text="🔒 Add Password", command=self.add_password_protection).pack(side=LEFT, padx=5)
        Button(top_frame2, text="❌ Remove Password", command=self.remove_password_protection).pack(side=LEFT, padx=5)
        Button(top_frame1, text="🔖 Add Bookmark", command=self.add_bookmark).pack(side=LEFT, padx=5)
        Button(top_frame1, text="❌ Remove Bookmark", command=self.remove_bookmark).pack(side=LEFT, padx=5)
        Button(top_frame2, text="Change Password", command=self.change_password).pack(side=LEFT, padx=5)
        
        self.eraser_button = Button(top_frame2, text="Eraser 🧽", command=self.toggle_eraser_mode)

        self.eraser_button.pack(side=LEFT, padx=5)
        self.bookmark_listbox = Listbox(top_frame1, height=3, width=15)
        self.bookmark_listbox.pack(side=LEFT, padx=5)
        self.bookmark_listbox.bind("<<ListboxSelect>>", self.goto_bookmark)

        self.page_count_label = Label(top_frame2, text="Total Pages: 0", fg="white", bg="#1e1e1e")
        self.page_count_label.pack(side=LEFT, padx=10)

        self.current_page_label = Label(top_frame2, text="Page 0 / 0", fg="white", bg="#1e1e1e", font=("Segoe UI", 11))
        self.current_page_label.pack(side=RIGHT, padx=10)




        self.canvas = Canvas(self.root, bg="#1e1e1e", highlightthickness=0)
        self.canvas.pack(side=LEFT, fill=BOTH, expand=True)

        scrollbar = Scrollbar(self.root, orient="vertical", command=self.canvas.yview)
        scrollbar.pack(side=RIGHT, fill=Y)
        self.canvas.configure(yscrollcommand=scrollbar.set)

        self.inner_frame = Frame(self.canvas, bg="#1e1e1e")
        self.canvas_window = self.canvas.create_window((0, 0), window=self.inner_frame, anchor="nw")
        

        self.canvas.bind("<Configure>", self.center_frame)
        self.canvas.bind("<Enter>", self._bind_scroll)
        self.canvas.bind("<Leave>", self._unbind_scroll)
     
        self.canvas.bind("<MouseWheel>", self.update_current_page_on_scroll)
        self.canvas.bind("<Button-4>", self.update_current_page_on_scroll)
        self.canvas.bind("<Button-5>", self.update_current_page_on_scroll)
        self.canvas.bind("<ButtonRelease-1>", self.update_current_page_on_scroll)

        self.canvas.bind_all("<ButtonRelease-1>", self.update_current_page_on_scroll)
        self.canvas.bind_all("<Motion>", self.update_current_page_on_scroll)
        
        self.canvas.bind("<Button-4>", lambda event: self.root.after(50, self.update_current_page_on_scroll))
        self.canvas.bind("<Button-5>", lambda event: self.root.after(50, self.update_current_page_on_scroll))
        self.canvas.bind("<ButtonRelease-1>", lambda event: self.root.after(50, self.update_current_page_on_scroll))

    

    def update_current_page_on_scroll(self, event=None):
        if not self.page_widgets:
            return

        scroll_y = self.canvas.canvasy(0)
        min_diff = float("inf")
        closest_index = 0

        for i, widget in enumerate(self.page_widgets):
            try:
                y = widget.winfo_y()
                diff = abs(y - scroll_y)
                if diff < min_diff:
                    min_diff = diff
                    closest_index = i
            except:
                continue

        self.current_page = closest_index

        if hasattr(self, "current_page_label"):
            self.current_page_label.config(text=f"Page {self.current_page + 1 } / {self.total_pages}")



    def delayed_page_update(self, event=None):
        self.root.after_idle(self.update_current_page_on_scroll)


    



    


    def center_frame(self, event):
        """Center the inner frame in the canvas."""
        canvas_width = event.width
        self.canvas.itemconfig(self.canvas_window, width=canvas_width)

    def _bind_scroll(self, event):
        """Bind mouse wheel for scrolling."""
        self.canvas.bind_all("<MouseWheel>", self._on_mousewheel)

    def _unbind_scroll(self, event):
        """Unbind mouse wheel when leaving canvas."""
        self.canvas.unbind_all("<MouseWheel>")

    def _on_mousewheel(self, event):
        """Handle mouse wheel scrolling."""
        self.canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")

    def toggle_draw_mode(self):
        """Toggle drawing mode for PDF annotations."""
        self.drawing = not self.drawing
        btn = None
        for widget in self.root.winfo_children():
            if isinstance(widget, Frame):
                for child in widget.winfo_children():
                    if isinstance(child, Button) and "Draw Mode" in child.cget("text"):
                        btn = child
                        break
        if btn:
            btn.config(bg="green" if self.drawing else "SystemButtonFace", fg="white" if self.drawing else "black")

    def start_draw(self, event):
        """Start drawing on canvas."""
        if self.drawing:
            event.widget.old_coords = (event.x, event.y)

    def draw(self, event):
        if not hasattr(self, "draw_highlights"):
            self.draw_highlights = {}

        canvas = event.widget
        wrapper = canvas.master
        try:
            current_page = self.page_widgets.index(wrapper)
        except ValueError:
            return

        # ✅ Eraser mode — doesn't need old_coords
        if self.eraser_mode:
            x, y = event.x, event.y
            lines = canvas.find_all()
            for item in lines:
                coords = canvas.coords(item)
                if len(coords) >= 4:
                    x0, y0, x1, y1 = coords[:4]
                    if abs(x - x0) < 10 and abs(y - y0) < 10:
                        canvas.delete(item)
            return  # Don't do any drawing

        # ✅ If drawing mode is OFF, do nothing
        if not self.drawing:
            return

        # ✅ Drawing mode — needs old_coords
        if not hasattr(canvas, "old_coords"):
            return

        x1, y1 = canvas.old_coords
        x2, y2 = event.x, event.y
        canvas.old_coords = (x2, y2)

        canvas.create_line(x1, y1, x2, y2, fill="red", width=2)
        self.draw_highlights.setdefault(current_page, []).append(
            (min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2))
        )


    def toggle_eraser_mode(self):
        self.eraser_mode = not self.eraser_mode

        if self.eraser_mode:
            self.eraser_button.config(bg="green", fg="white")
            self.canvas.config(cursor="dotbox")  # ✅ Eraser-style cursor
            messagebox.showinfo("Eraser", "Eraser mode ON")
        else:
            self.eraser_button.config(bg="SystemButtonFace", fg="black")
            self.canvas.config(cursor="arrow")  # ✅ Normal cursor
            messagebox.showinfo("Eraser", "Eraser mode OFF")



    def stop_draw(self, event):
        """Stop drawing on canvas."""
        if hasattr(event.widget, "old_coords"):
            del event.widget.old_coords

    def get_total_pages(self, file_path, ext):
        """Calculate the total number of pages for the given file."""
        file_path = os.path.normpath(file_path)
        try:
            if ext == ".pdf":
                doc = fitz.open(file_path)
                if doc.is_encrypted:
                    password = self.prompt_password(file_path)
                    if password is None:
                        doc.close()
                        raise Exception("No password provided.")
                    if not doc.authenticate(password):
                        doc.close()
                        raise Exception("Incorrect password.")
                total_pages = len(doc)
                doc.close()
                return total_pages
            elif ext == ".docx":
                doc = Document(file_path)
                lines = sum(1 for para in doc.paragraphs if para.text.strip())
                return max(1, lines // 30)
            elif ext == ".pptx":
                ppt = Presentation(file_path)
                return len(ppt.slides)
            elif ext == ".sps":
                if file_path in self.protected_files:
                    password = self.prompt_password(file_path)
                    if password is None:
                        raise ValueError("No password provided.")
                    kdf = PBKDF2HMAC(
                        algorithm=hashes.SHA256(),
                        length=32,
                        salt=b'salt_',
                        iterations=100000,
                    )
                    key = base64.urlsafe_b64encode(kdf.derive(password.encode('utf-8')))
                    fernet = Fernet(key)
                    with open(file_path, "rb") as f:
                        encrypted_data = f.read()
                    try:
                        decrypted_data = fernet.decrypt(encrypted_data)
                        lines = decrypted_data.decode('utf-8').splitlines()
                        return max(1, len(lines) // 30)
                    except Exception:
                        raise ValueError("Failed to decrypt SPS file.")
                else:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        lines = f.readlines()
                    return max(1, len(lines) // 30)
            elif ext == ".epub":
                book = epub.read_epub(file_path)
                return sum(1 for item in book.get_items() if item.get_type() == ITEM_DOCUMENT)
            else:
                return 1
        except Exception as e:
            logging.error(f"Failed to calculate pages for {file_path}: {str(e)}")
            messagebox.showerror("Error", f"Failed to calculate pages: {str(e)}")
            return 1

    
    def select_file(self):
        file_path = filedialog.askopenfilename(
            filetypes=[("Supported Files", ".pdf *.docx *.pptx *.sps *.epub"), ("All Files", ".*")]
        )
        if not file_path:
            return

        # Close existing pdf_doc
        if self.pdf_doc:
            self.pdf_doc.close()
            self.pdf_doc = None

        self.file_path = os.path.normpath(file_path)
        ext = os.path.splitext(self.file_path)[1].lower()

        try:
            # Calculate total pages
            total_pages = self.get_total_pages(self.file_path, ext)
            self.total_pages = total_pages
            self.current_page = 0  # Reset current page
            self.page_count_label.config(text=f"Total Pages: {total_pages}")
            self.current_page_label.config(text=f"Page 1 / {total_pages}")  # Initialize label

            if ext == ".pdf":
                doc = fitz.open(self.file_path)
                if doc.is_encrypted:
                    password = simpledialog.askstring("Password", f"Enter password for {os.path.basename(self.file_path)}:", show="*")
                    if password is None:
                        doc.close()
                        messagebox.showerror("Error", "No password provided.")
                        return
                    if not doc.authenticate(password):
                        doc.close()
                        messagebox.showerror("Error", "Incorrect password.")
                        return
                    messagebox.showinfo("Success", "Password verified successfully.")

                self.draw_highlights = {}
                self.search_highlights = {}
                self.pdf_doc = doc
                self.doc_text = "".join([page.get_text() for page in doc])
                self.display_pdf(doc)
                self.update_current_page_on_scroll()  # Update page display after loading

            elif ext == ".docx":
                self.doc_text = self.extract_docx(self.file_path)
                self.display_docx(self.file_path)
                self.update_current_page_on_scroll()  # Update for non-PDF files

            elif ext == ".pptx":
                self.doc_text = self.extract_pptx(self.file_path)
                self.display_pptx(self.file_path)
                self.update_current_page_on_scroll()

            elif ext == ".sps":
                self.doc_text = self.extract_text(self.file_path)
                self.display_text(self.doc_text)
                self.update_current_page_on_scroll()

            elif ext == ".epub":
                self.doc_text = self.extract_epub(self.file_path)
                self.display_text(self.doc_text)
                self.update_current_page_on_scroll()

            else:
                messagebox.showerror("Error", "Unsupported file format.")

        except Exception as e:
            logging.error(f"Failed to open file {self.file_path}: {str(e)}")
            messagebox.showerror("Error", f"Failed to open file: {str(e)}")
            if self.pdf_doc:
                self.pdf_doc.close()
                self.pdf_doc = None

    def extract_docx(self, path):
        from docx import Document
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


    def extract_pptx(self, path):
        """Extract text from a PPTX file."""
        try:
            ppt = Presentation(path)
            content = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
            return content
        except Exception as e:
            raise ValueError(f"Failed to open PPTX: {str(e)}")

    def extract_text(self, path):
        """Extract text from a text file, decrypting if necessary."""
        path = os.path.normpath(path)
        try:
            if path in self.protected_files:
                password = self.prompt_password(path)
                if password is None:
                    raise ValueError("No password provided.")
                kdf = PBKDF2HMAC(
                    algorithm=hashes.SHA256(),
                    length=32,
                    salt=b'salt_',
                    iterations=100000,
                )
                key = base64.urlsafe_b64encode(kdf.derive(password.encode('utf-8')))
                fernet = Fernet(key)
                with open(path, "rb") as f:
                    encrypted_data = f.read()
                try:
                    decrypted_data = fernet.decrypt(encrypted_data)
                    return decrypted_data.decode('utf-8')
                except Exception:
                    raise ValueError("Failed to decrypt text file: Incorrect password or corrupted file.")
            else:
                with open(path, "r", encoding="utf-8") as f:
                    return f.read()
        except Exception as e:
            raise ValueError(f"Failed to open text file: {str(e)}")

    def extract_epub(self, path):
        """Extract text from an EPUB file."""
        try:
            book = epub.read_epub(path)
            content = ""
            for item in book.get_items():
                if item.get_type() == ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    content += html2text.html2text(soup.get_text())
            return content
        except Exception as e:
            raise ValueError(f"Failed to open EPUB: {str(e)}")

    def display_text(self, text):
        """Display text content in the UI."""
        for widget in self.inner_frame.winfo_children():
            widget.destroy()

        lines = text.strip().split('\n')
        lines_per_page = 30
        pages = [lines[i:i+lines_per_page] for i in range(0, len(lines), lines_per_page)]
        self.page_widgets.clear()

        for page_lines in pages:
            wrapper = Frame(self.inner_frame, bg="#1e1e1e", pady=20)
            wrapper.pack()

            text_widget = Text(wrapper, wrap=WORD, font=("Segoe UI", 12), bg="white", height=30)
            text_widget.pack(fill=BOTH, expand=True)
            text_widget.insert(END, "\n".join(page_lines).strip())
            text_widget.config(state=DISABLED)
            self.current_text_widget = text_widget

            self.page_widgets.append(wrapper)

        self.inner_frame.update_idletasks()
        self.canvas.config(scrollregion=self.canvas.bbox("all"))
        self.current_text_widget = text_widget if self.page_widgets else None

    def display_docx(self, path):
        """Display DOCX content as HTML."""
        try:
            with open(path, "rb") as docx_file:
                result = mammoth.convert_to_html(docx_file)
                html_content = result.value
                messages = result.messages
                if messages:
                    critical_messages = [m for m in messages if "error" in str(m).lower()]
                    if critical_messages:
                        messagebox.showwarning("Conversion Warnings", "\n".join(str(m) for m in critical_messages))
            self.display_html(html_content)
        except Exception as e:
            messagebox.showerror("Error", f"Failed to display DOCX: {str(e)}")

    def display_pptx(self, path):
        """Display PPTX content as HTML with enhanced formatting."""
        try:
            ppt = Presentation(path)
            html_content = "<html><body><h1>Presentation Slides</h1>"
            for i, slide in enumerate(ppt.slides, 1):
                html_content += f"<h2>Slide {i}</h2>"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.replace("\n", "<br>")
                        style = ""
                        if shape.has_text_frame:
                            paragraph = shape.text_frame.paragraphs[0]
                            font_size = paragraph.font.size.pt if paragraph.font.size else 12
                            is_bold = "bold" if paragraph.font.bold else "normal"
                            color = paragraph.font.color.rgb if paragraph.font.color.type == 1 else None
                            align = paragraph.alignment if paragraph.alignment else None
                            style_attrs = [f"font-size: {font_size}px", f"font-weight: {is_bold}"]
                            if color:
                                color_hex = f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"
                                style_attrs.append(f"color: {color_hex}")
                            if align:
                                align_map = {0: "left", 1: "center", 2: "right", 3: "justify"}
                                style_attrs.append(f"text-align: {align_map.get(align, 'left')}")
                            style = f"style='{'; '.join(style_attrs)}'"
                        html_content += f"<p {style}>{text}</p>"
            html_content += "</body></html>"
            self.display_html(html_content)
        except Exception as e:
            messagebox.showerror("Error", f"Failed to display PPTX: {str(e)}")

    def display_html(self, html_content):
        """Display HTML content in the UI."""
        for widget in self.inner_frame.winfo_children():
            widget.destroy()

        self.page_widgets.clear()
        wrapper = Frame(self.inner_frame, bg="#1e1e1e", pady=20)
        wrapper.pack(fill=BOTH, expand=True)

        html_widget = HTMLLabel(wrapper, html=html_content, background="white")
        html_widget.pack(fill=BOTH, expand=True)

        self.page_widgets.append(wrapper)
        self.inner_frame.update_idletasks()
        self.canvas.config(scrollregion=self.canvas.bbox("all"))
        self.current_text_widget = None

    def display_pdf(self, doc, highlights=None, dpi=120):
        """Display PDF pages using PyMuPDF with optional highlights."""
        self.total_pages = len(doc)
        self.tk_images.clear()
        self.page_widgets.clear()
        for widget in self.inner_frame.winfo_children():
            widget.destroy()

        try:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=dpi)
                pil_img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGBA")

                if highlights and i in highlights:
                    draw = ImageDraw.Draw(pil_img, "RGBA")
                    scale = dpi / 72
                    for rect in highlights[i]:
                        x0, y0, x1, y1 = rect
                        box = [x0 * scale, y0 * scale, x1 * scale, y1 * scale]
                        draw.rectangle([box[0], box[3] - 4, box[2], box[3]], fill=(255, 255, 0, 180))

                tk_img = ImageTk.PhotoImage(pil_img)
                self.tk_images.append(tk_img)

                wrapper = Frame(self.inner_frame, bg="#1e1e1e", pady=20)
                wrapper.pack()
                canvas = Canvas(wrapper, width=pil_img.width, height=pil_img.height, bg="white")
                canvas.pack()
                canvas.create_image(0, 0, anchor="nw", image=tk_img)

                canvas.bind("<ButtonPress-1>", self.start_draw)
                canvas.bind("<B1-Motion>", self.draw)
                canvas.bind("<ButtonRelease-1>", self.stop_draw)

                self.page_widgets.append(wrapper)

            self.inner_frame.update_idletasks()
            self.canvas.config(scrollregion=self.canvas.bbox("all"))
            self.current_text_widget = None
            self.current_page = 0  # Reset to first page
            self.current_page_label.config(text=f"Page 1 / {self.total_pages}")  # Update label
            self.update_current_page_on_scroll()  # Ensure initial page is set

        except Exception as e:
            logging.error(f"Failed to display PDF: {str(e)}")
            messagebox.showerror("PDF Error", f"❌ Failed to display PDF:\n{e}")




    def search_keyword(self):
        """Search for a keyword and highlight matches."""
        keyword = self.text_entry.get().strip().lower()
        if not keyword:
            return

        found = False
        if self.file_path.endswith(".pdf"):
            if not self.pdf_doc:
                messagebox.showerror("Error", "PDF not loaded.")
                return
            try:
                if self.pdf_doc.is_encrypted:
                    password = self.prompt_password(self.file_path)
                    if password is None:
                        messagebox.showerror("PDF Error", "No password provided.")
                        return
                    if not self.pdf_doc.authenticate(password):
                        messagebox.showerror("PDF Error", "Incorrect password.")
                        return
                    logging.debug(f"Password verified with PyMuPDF for {self.file_path}")
                results = {}
                for i, page in enumerate(self.pdf_doc):
                    words = page.get_text("words")
                    matches = [fitz.Rect(w[0], w[1], w[2], w[3]) for w in words if w[4].lower() == keyword]
                    if matches:
                        results[i] = matches
                        found = True
                self.display_pdf(self.pdf_doc, highlights=results if results else None)
            except Exception as e:
                logging.error(f"Failed to search PDF {self.file_path}: {str(e)}")
                messagebox.showerror("PDF Error", f"❌ Failed to search PDF:\n{e}")
        elif self.current_text_widget:
            self.current_text_widget.tag_remove("highlight", "1.0", END)
            idx = "1.0"
            while True:
                idx = self.current_text_widget.search(keyword, idx, nocase=1, stopindex=END)
                if not idx:
                    break
                lastidx = f"{idx}+{len(keyword)}c"
                self.current_text_widget.tag_add("highlight", idx, lastidx)
                idx = lastidx
                found = True
            self.current_text_widget.tag_config("highlight", background="yellow", foreground="black")

        if not found:
            messagebox.showinfo("Search Result", f"No matches found for '{keyword}'.")

    def clear_highlights(self):
        """Clear search highlights."""
        if self.file_path.endswith(".pdf"):
            if self.pdf_doc:
                try:
                    if self.pdf_doc.is_encrypted:
                        password = self.prompt_password(self.file_path)
                        if password is None:
                            messagebox.showerror("Error", "No password provided.")
                            return
                        if not self.pdf_doc.authenticate(password):
                            messagebox.showerror("Error", "Incorrect password.")
                            return
                        logging.debug(f"Password verified with PyMuPDF for {self.file_path}")
                    self.display_pdf(self.pdf_doc)
                except Exception as e:
                    logging.error(f"Failed to clear highlights for PDF: {str(e)}")
                    messagebox.showerror("PDF Error", f"❌ Failed to clear highlights:\n{e}")
            else:
                messagebox.showerror("Error", "PDF not loaded.")
        elif self.current_text_widget:
            self.current_text_widget.tag_remove("highlight", "1.0", END)

    def copy_text_popup(self):
        """Display extracted text in a popup with scrollbar."""
        if not self.doc_text:
            messagebox.showinfo("Info", "No text available to copy.")
            return

        popup = Toplevel(self.root)
        popup.title("📋 Document Text Content")
        popup.geometry("800x600")

        frame = Frame(popup)
        frame.pack(fill=BOTH, expand=True)

        scrollbar = Scrollbar(frame)
        scrollbar.pack(side=RIGHT, fill=Y)

        text_box = Text(frame, wrap=WORD, yscrollcommand=scrollbar.set)
        text_box.pack(side=LEFT, fill=BOTH, expand=True)
        scrollbar.config(command=text_box.yview)

        text_box.insert(END, self.doc_text.strip())
        text_box.config(state=DISABLED)

    def summarize_text(self):
        """Fast extractive summary for large PDFs."""
        text = self.doc_text.strip()
        if len(text) < 100:
            messagebox.showinfo("Too Short", "Text too short to summarize.")
            return

        processing_popup = Toplevel(self.root)
        processing_popup.title("Processing")
        Label(processing_popup, text="Summarizing... Please wait.", font=("Segoe UI", 12)).pack(pady=20, padx=20)
        processing_popup.update()

        try:
            parser = PlaintextParser.from_string(text, Tokenizer("english"))
            summarizer = LexRankSummarizer()

            # Select number of sentences based on document length
            total_len = len(self.doc_text)
            if total_len < 3000:
                sentence_count = 5
            elif total_len < 12000:
                sentence_count = 10
            else:
                sentence_count = 20

            summary_sentences = summarizer(parser.document, sentence_count)
            final_summary = " ".join(str(s) for s in summary_sentences)

            processing_popup.destroy()
            self.show_summary(final_summary)

        except Exception as e:
            processing_popup.destroy()
            self.show_summary(f"⚠ Error: {e}")




    def show_summary(self, summary):
        """Display summary in a popup."""
        popup = Toplevel(self.root)
        popup.title("🧠 Summary")
        popup.geometry("800x400")

        text_box = Text(popup, wrap=WORD, font=("Segoe UI", 12))
        text_box.pack(fill=BOTH, expand=True)
        text_box.insert(END, summary)
        text_box.config(state=DISABLED)

    def add_bookmark(self):
        """Add a bookmark using scroll position directly."""

    # Scroll ka y position
        y = self.canvas.canvasy(0)
        y = round(y, 2)

    # Total height and page calculation
        total_height = self.inner_frame.winfo_height()
        page_height = total_height / self.total_pages if self.total_pages else 1
        page_number = int(y // page_height) + 1

    # Duplicate check on page_number
        if any(b["page"] == page_number for b in self.bookmarks):
            messagebox.showinfo("🔖 Bookmark", f"Page {page_number} is already bookmarked.")
            return

    # Save the bookmark
        self.bookmarks.append({"y": y, "page": page_number})
        self.refresh_bookmark_listbox()









    def remove_bookmark(self):
        selection = self.bookmark_listbox.curselection()
        if not selection:
            return
        index = selection[0]
        self.bookmarks.pop(index)
        self.refresh_bookmark_listbox()




    def goto_bookmark(self, event):
        selection = self.bookmark_listbox.curselection()
        if not selection:
            return
        bookmark = self.bookmarks[selection[0]]
        y = bookmark["y"]

        total_height = self.inner_frame.winfo_height()
        if total_height > 0:
            self.canvas.yview_moveto(y / total_height)



    def refresh_bookmark_listbox(self):
        self.bookmark_listbox.delete(0, END)
        for bookmark in self.bookmarks:
            page = bookmark["page"]
            self.bookmark_listbox.insert(END, f"Page {page}")



    def next_page(self):
        if self.current_page < self.total_pages - 1:
            self.show_page(self.current_page + 1)

    def prev_page(self):
        if self.current_page > 0:
            self.show_page(self.current_page - 1)


    def show_page(self, page_number):
        self.current_page = page_number
    # Show page normally

   
    def get_total_pages(self, file_path, ext):
        """Calculate the total number of pages for the given file."""
        file_path = os.path.normpath(file_path)
        try:
            if ext == ".pdf":
                doc = fitz.open(file_path)
                if doc.is_encrypted:
                    password = self.prompt_password(file_path)
                    if password is None:
                        doc.close()
                        raise Exception("No password provided.")
                    if not doc.authenticate(password):
                        doc.close()
                        raise Exception("Incorrect password.")
                total_pages = len(doc)
                doc.close()
                return total_pages
            elif ext == ".docx":
                doc = Document(file_path)
                lines = sum(1 for para in doc.paragraphs if para.text.strip())
                return max(1, lines // 30)
            elif ext == ".pptx":
                ppt = Presentation(file_path)
                return len(ppt.slides)
            elif ext == ".sps":
                if file_path in self.protected_files:
                    password = self.prompt_password(file_path)
                    if password is None:
                        raise ValueError("No password provided.")
                    kdf = PBKDF2HMAC(
                        algorithm=hashes.SHA256(),
                        length=32,
                        salt=b'salt_',
                        iterations=100000,
                    )
                    key = base64.urlsafe_b64encode(kdf.derive(password.encode('utf-8')))
                    fernet = Fernet(key)
                    with open(file_path, "rb") as f:
                        encrypted_data = f.read()
                    try:
                        decrypted_data = fernet.decrypt(encrypted_data)
                        lines = decrypted_data.decode('utf-8').splitlines()
                        return max(1, len(lines) // 30)
                    except Exception:
                        raise ValueError("Failed to decrypt SPS file.")
                else:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        lines = f.readlines()
                    return max(1, len(lines) // 30)
            elif ext == ".epub":
                book = epub.read_epub(file_path)
                return sum(1 for item in book.get_items() if item.get_type() == ITEM_DOCUMENT)
            else:
                return 1
        except Exception as e:
            logging.error(f"Failed to calculate pages for {file_path}: {str(e)}")
            messagebox.showerror("Error", f"Failed to calculate pages: {str(e)}")
            return 1



    def add_password_protection(self):
        """Add password protection to a selected file."""
        file_path = filedialog.askopenfilename(
            filetypes=[("Supported Files", ".pdf *.docx *.pptx *.sps *.epub"), ("All Files", ".*")]
        )
        if not file_path:
            return

        file_path = os.path.normpath(file_path)
        if file_path in self.protected_files:
            messagebox.showwarning("Warning", "This file is already password-protected.")
            return

        password = simpledialog.askstring("Set Password", f"Enter password for {os.path.basename(file_path)}:", show="*")
        if not password or not password.strip():
            messagebox.showerror("Error", "Password cannot be empty.")
            return

        confirm_password = simpledialog.askstring("Confirm Password", "Confirm password:", show="*")
        if password != confirm_password:
            messagebox.showerror("Error", "Passwords do not match.")
            return

        password = password.strip()
        logging.debug(f"Setting password for {file_path} (length: {len(password)})")

        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            try:
                pypdf2_version = pkg_resources.get_distribution("PyPDF2").version
                if pypdf2_version < "1.26.0":
                    messagebox.showwarning(
                        "Warning",
                        f"Your PyPDF2 version ({pypdf2_version}) is outdated. "
                        "Please upgrade to version 1.26.0 or higher for AES-128 encryption.\n"
                        "Run: pip install --upgrade PyPDF2"
                    )

                reader = PdfReader(file_path)
                if reader.is_encrypted:
                    messagebox.showerror("Error", "This PDF is already encrypted.")
                    return

                output_path = os.path.normpath(os.path.splitext(file_path)[0] + "_protected.pdf")
                protect_pdf(file_path, output_path, password)

                # Verify encryption with PyMuPDF (since we'll open it with fitz)
                if not self.test_pdf_encryption_fitz(output_path, password):
                    messagebox.showerror("Error", "Failed to verify encrypted PDF. Encryption may be corrupted.")
                    logging.error(f"Encryption verification failed for {output_path}")
                    return

                self.protected_files[output_path] = password
                self.save_protected_files()
                messagebox.showinfo(
                    "Success",
                    f"Password-protected PDF created: {os.path.basename(output_path)}.\n"
                    "This file will require a password in other PDF readers."
                )
            except Exception as e:
                messagebox.showerror("Error", f"Failed to encrypt PDF: {str(e)}")
                logging.error(f"Failed to encrypt PDF {file_path}: {str(e)}")
        elif ext == ".sps":
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    text_content = f.read()
                kdf = PBKDF2HMAC(
                    algorithm=hashes.SHA256(),
                    length=32,
                    salt=b'salt_',
                    iterations=100000,
                )
                key = base64.urlsafe_b64encode(kdf.derive(password.encode('utf-8')))
                fernet = Fernet(key)
                encrypted_content = fernet.encrypt(text_content.encode('utf-8'))
                output_path = os.path.normpath(os.path.splitext(file_path)[0] + "_protected.sps")
                with open(output_path, "wb") as f:
                    f.write(encrypted_content)
                self.protected_files[output_path] = self.hash_password(password)
                self.save_protected_files()
                messagebox.showinfo(
                    "Success",
                    f"Encrypted text file created: {os.path.basename(output_path)}.\n"
                    "This file can only be opened in SmartDocReader with the correct password."
                )
            except Exception as e:
                messagebox.showerror("Error", f"Failed to encrypt text file: {str(e)}")
                logging.error(f"Failed to encrypt SPS {file_path}: {str(e)}")
        else:
            self.protected_files[file_path] = self.hash_password(password)
            self.save_protected_files()
            messagebox.showinfo(
                "Success",
                f"Password protection added for {os.path.basename(file_path)}.\n"
                "Note: This protection applies only within SmartDocReader for non-PDF/SPS files."
            )
    def remove_password_protection(self):
        """Remove password protection from the currently opened file."""
        if not self.file_path:
            messagebox.showwarning("Warning", "No file is currently open.")
            return

        file_path = self.file_path
        ext = os.path.splitext(file_path)[1].lower()

        # Check if the file is password-protected
        if ext == ".pdf":
            try:
                reader = PdfReader(file_path)
                if not reader.is_encrypted and file_path not in self.protected_files:
                    messagebox.showinfo("Info", "This PDF is not password-protected.")
                    return
            except Exception as e:
                logging.error(f"Error reading PDF encryption status: {str(e)}")
                messagebox.showerror("Error", f"Error reading file: {str(e)}")
                return
        else:
            if file_path not in self.protected_files:
                messagebox.showinfo("Info", "This file is not password-protected.")
                return

        # Ask for the current password
        password = simpledialog.askstring("Password", f"Enter current password for {os.path.basename(file_path)}:", show="*")
        if not password:
            return

        try:
            if ext == ".pdf":
                # Verify and decrypt PDF
                reader = PdfReader(file_path)
                if reader.is_encrypted:
                    if reader.decrypt(password) == 0:
                        messagebox.showerror("Error", "Incorrect password.")
                        return
                else:
                    messagebox.showinfo("Info", "This PDF is not encrypted.")
                    return

                # Verify password for fitz (for self.pdf_doc compatibility)
                if not self.test_pdf_encryption_fitz(file_path, password):
                    messagebox.showerror("Error", "Incorrect password for PDF rendering.")
                    return

                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)

                # Save decrypted PDF
                output_path = file_path
                with open(output_path, "wb") as f:
                    writer.write(f)

                del self.protected_files[file_path]
                self.save_protected_files()
                messagebox.showinfo("Success", "Password removed from PDF.")

                # Reload self.pdf_doc if the current file is affected
                if self.pdf_doc and self.file_path == file_path:
                    self.pdf_doc.close()
                    self.pdf_doc = fitz.open(file_path)
                    self.display_pdf(self.pdf_doc)

            elif ext == ".sps":
                # Decrypt SPS file
                kdf = PBKDF2HMAC(algorithm=hashes.SHA256(), length=32, salt=b'salt_', iterations=100000)
                key = base64.urlsafe_b64encode(kdf.derive(password.encode('utf-8')))
                fernet = Fernet(key)

                with open(file_path, "rb") as f:
                    decrypted_data = fernet.decrypt(f.read())

                # Overwrite file with decrypted text
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(decrypted_data.decode('utf-8'))

                del self.protected_files[file_path]
                self.save_protected_files()
                messagebox.showinfo("Success", "Password removed from SPS file.")

            else:
                # Non-encrypted files: just remove password from records
                del self.protected_files[file_path]
                self.save_protected_files()
                messagebox.showinfo("Success", "Password removed from file.")

        except Exception as e:
            logging.error(f"Failed to remove password from {file_path}: {str(e)}")
            messagebox.showerror("Error", f"Failed to remove password: {str(e)}")


    def change_password(self):
    # ✅ Don't close the PDF at start, first check if password is needed
        if not self.file_path:
            messagebox.showerror("Error", "No file is currently open.")
            return

        ext = os.path.splitext(self.file_path)[1].lower()
        old_password = simpledialog.askstring("Current Password", "Enter current password:", show="*")
        if not old_password:
            return

        try:
            if ext == ".pdf":
                reader = PdfReader(self.file_path)
                if not reader.is_encrypted:
                # ✅ If not encrypted, don't close the PDF → show info and return
                    messagebox.showinfo("Info", "This PDF is not password protected.")
                    return

                if reader.decrypt(old_password) == 0:
                    messagebox.showerror("Error", "Incorrect current password.")
                    return

            # ✅ NOW close the open file before changing it
                if hasattr(self, "pdf_doc") and self.pdf_doc:
                    try:
                        self.pdf_doc.close()
                        self.pdf_doc = None
                    except Exception as e:
                        print(f"Error closing PDF before password change: {e}")

            # ✅ New password input
                new_password = simpledialog.askstring("New Password", "Enter new password:", show="*")
                confirm_password = simpledialog.askstring("Confirm Password", "Confirm new password:", show="*")

                if new_password != confirm_password:
                    messagebox.showerror("Error", "Passwords do not match.")
                    return

                writer = PdfWriter()
                for page in reader.pages:
                    writer.add_page(page)

                writer.encrypt(user_password=new_password, owner_password=new_password)

                temp_output = self.file_path + ".temp"
                with open(temp_output, "wb") as temp_f:
                    writer.write(temp_f)

                os.remove(self.file_path)
                os.rename(temp_output, self.file_path)

                self.protected_files[self.file_path] = new_password
                self.save_protected_files()

                messagebox.showinfo("Success", "Password changed successfully.")

            # ✅ Now reopen the modified PDF
                self.reopen_pdf_with_new_password(new_password)

            else:
                messagebox.showerror("Error", "This version only supports PDF password changes.")

        except Exception as e:
            messagebox.showerror("Error", f"Failed to change password:\n{str(e)}")



    def reopen_pdf_with_new_password(self, password):
        try:
            doc = fitz.open(self.file_path)
            if doc.is_encrypted:
                if not doc.authenticate(password):
                    messagebox.showerror("Error", "Failed to open PDF with new password.")
                    return

            self.pdf_doc = doc
            self.doc_text = "".join([page.get_text() for page in doc])

        # Optionally redisplay it
            self.display_pdf(doc)

        except Exception as e:
            messagebox.showerror("Error", f"Error reopening PDF: {str(e)}")


if __name__ == "__main__":
    root = Tk()
    app = SmartDocReader(root)
    root.mainloop()

